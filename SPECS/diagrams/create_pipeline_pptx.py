"""Generate PowerPoint with 2 slides comparing RAG pipelines side by side"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

base_path = "/home/tromanow/COHORT/TX/SPECS/diagrams"

# Create presentation (widescreen 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# Slide 1: Custom RAG vs Kendra
slide1 = prs.slides.add_slide(blank_layout)
title1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.4))
title1.text_frame.paragraphs[0].text = "Custom RAG vs Kendra"
title1.text_frame.paragraphs[0].font.size = Pt(24)
title1.text_frame.paragraphs[0].font.bold = True
title1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

slide1.shapes.add_picture(os.path.join(base_path, "custom-rag-1.png"), Inches(1.0), Inches(0.6), width=Inches(3.9))
slide1.shapes.add_picture(os.path.join(base_path, "kendra-1.png"), Inches(7.5), Inches(0.6), width=Inches(3.0))

# Slide 2: OpenAI vs Vertex AI
slide2 = prs.slides.add_slide(blank_layout)
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.4))
title2.text_frame.paragraphs[0].text = "OpenAI vs Vertex AI"
title2.text_frame.paragraphs[0].font.size = Pt(24)
title2.text_frame.paragraphs[0].font.bold = True
title2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

slide2.shapes.add_picture(os.path.join(base_path, "openai-agent-1.png"), Inches(0.2), Inches(0.6), width=Inches(6.4))
slide2.shapes.add_picture(os.path.join(base_path, "vertex-agent-1.png"), Inches(6.7), Inches(0.6), width=Inches(6.4))

# Save
output_path = os.path.join(base_path, "rag_pipelines_comparison.pptx")
prs.save(output_path)
print(f"Created: {output_path}")
